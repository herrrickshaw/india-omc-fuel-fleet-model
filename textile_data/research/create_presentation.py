#!/usr/bin/env python3
"""
Generate PowerPoint presentation: Textile & Chemical Sector Investment Analysis
Master Investment List + Import Substitution Opportunity Matrix
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import datetime

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
TITLE_COLOR = RGBColor(0, 51, 102)      # Dark blue
ACCENT_COLOR = RGBColor(255, 102, 0)   # Orange
TEXT_COLOR = RGBColor(51, 51, 51)       # Dark gray
LIGHT_BG = RGBColor(242, 242, 242)      # Light gray

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = TITLE_COLOR

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_COLOR

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.3))
    date_frame = date_box.text_frame
    p = date_frame.paragraphs[0]
    p.text = f"Investment Analysis | {datetime.date.today().strftime('%B %d, %Y')}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(200, 200, 200)

    return slide

def add_content_slide(prs, title, content_points=None, add_boxes=False):
    """Add content slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = TITLE_COLOR
    title_shape.line.color.rgb = TITLE_COLOR

    # Title text
    title_frame = title_shape.text_frame
    title_frame.margin_bottom = Inches(0.05)
    title_frame.margin_left = Inches(0.3)
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Content
    if content_points:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(6))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, point in enumerate(content_points):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = point
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
            p.level = 0
            p.space_before = Pt(6)
            p.space_after = Pt(6)

    return slide

def add_table_slide(prs, title, table_data):
    """Add slide with table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = TITLE_COLOR
    title_shape.line.color.rgb = TITLE_COLOR

    title_frame = title_shape.text_frame
    title_frame.margin_left = Inches(0.3)
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Add table
    rows, cols = len(table_data), len(table_data[0])
    left = Inches(0.3)
    top = Inches(1.2)
    width = Inches(9.4)
    height = Inches(5.8)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    for col_idx, col in enumerate(table_shape.columns):
        col.width = Inches(9.4 / cols)

    # Fill table data
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table_shape.cell(row_idx, col_idx)
            cell.text = str(cell_text)

            # Format header row
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TITLE_COLOR
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        run.font.bold = True
                        run.font.size = Pt(11)
            else:
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_BG
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(10)
                        run.font.color.rgb = TEXT_COLOR

    return slide

# ============================================================================
# SLIDE 1: TITLE SLIDE
# ============================================================================
add_title_slide(prs,
    "Textile & Chemical Sector Investment Analysis",
    "Master Investment List + Import Substitution Opportunity Matrix")

# ============================================================================
# SLIDE 2: EXECUTIVE SUMMARY
# ============================================================================
add_content_slide(prs, "Executive Summary", [
    "📊 Total Ecosystem Capex (FY26-FY30): ₹250,000+ Crore",
    "🎯 Annual Import Substitution Opportunity: ₹4,500-5,500 Crore",
    "🏆 Tier-1 Anchor Companies: 9 entities with ₹180,000 Cr (72% of capex)",
    "⚠️ CRITICAL DEADLINE: Mar 2026 for BPCL Bina LLDPE supply contracts",
    "📈 Margin Recovery Timeline: FY27-30 (INDUTECH +₹210-280M, CLOTHTECH +₹300-400M)",
    "🚀 New Company Targeting: 50-80 companies for investment in FY26-30 period"
])

# ============================================================================
# SLIDE 3: INVESTMENT LANDSCAPE OVERVIEW
# ============================================================================
add_content_slide(prs, "Investment Landscape (FY26-30)", [
    "Tier 1 (>₹5,000 Cr): ₹180,000 Cr — RIL O2C (₹75K), BPCL (₹18.9K), PM Mitra parks (₹70K)",
    "Tier 2 (₹1,000-5,000 Cr): ₹45,000 Cr — Trident, Jain Cord, Arvind, AB Cotspin, Indorama",
    "Tier 3 (₹300-1,000 Cr): ₹5,000 Cr — Atul, Bodal, Winsome, RSVPM, Best Corp",
    "Tier 4 (₹100-300 Cr): ₹3,000 Cr — Emerging/regional players, specialty niche companies",
    "Tier 5 (Recycling): ₹5,000 Cr — RE&UP (₹4.8K), Aquafil, circular economy players",
    "Government Support: ₹30,000 Cr — PLI (₹10.7K), PM Mitra, TEEM, Tex Eco, PCPIR"
])

# ============================================================================
# SLIDE 4: MASTER INVESTMENT LIST - TIER 1 COMPANIES
# ============================================================================
tier1_data = [
    ["Company", "Sector", "Capex (₹ Cr)", "Timeline", "Focus Area"],
    ["Reliance Industries", "Petrochemical", "75,000", "FY28-30", "+1.8 MMTPA PE, +1.0 PP"],
    ["BPCL (Bina + AP)", "Petrochemical", "18,900", "FY27-28", "+1.0 MMTPA LLDPE, Specialty"],
    ["PM Mitra Parks (7)", "Textile Mfg", "70,000", "FY26-30", "All segments, 300K+ jobs"],
    ["Trident Group", "Textiles", "4,881", "FY28", "HOMETECH (towels, rugs)"],
    ["Indorama Ventures", "Polyester", "3,000-5,000", "FY27-28", "+0.8 MMTPA Paradeep"]
]
add_table_slide(prs, "Tier-1 Companies (>₹5,000 Cr Capex)", tier1_data)

# ============================================================================
# SLIDE 5: TIER 2 COMPANIES
# ============================================================================
tier2_data = [
    ["Company", "Capex (₹ Cr)", "Segment", "Timeline", "Margin Opportunity"],
    ["Jain Cord Industries", "2,515", "INDUTECH", "FY28", "₹210-280M (LLDPE-dependent)"],
    ["AB Cotspin", "1,300", "Yarn supply", "FY27", "Ecosystem enabler"],
    ["Arvind Mills", "1,024", "CLOTHTECH", "FY28", "+₹300-400M margin swing"],
    ["Sanathan Textiles", "1,000", "Specialty polyester", "FY28", "Premium blends"],
    ["Best Corp Tiruppur", "832", "MOBILTECH", "FY28", "₹150-200M (specialty polymers)"]
]
add_table_slide(prs, "Tier-2 Companies (₹1,000-5,000 Cr)", tier2_data)

# ============================================================================
# SLIDE 6: IMPORT SUBSTITUTION OPPORTUNITY - OVERVIEW
# ============================================================================
add_content_slide(prs, "Import Substitution Landscape", [
    "🔴 PRIORITY 1 - LLDPE (Linear Low-Density PE): ₹1.8-2.0B import → 80% substitution by FY28",
    "    Current: 100% imported from Saudi Arabia | Supplier: BPCL Bina (+1.0 MMTPA)",
    "    Cost Savings: $150-200/MT | Affected Segment: INDUTECH nonwovens, geotextiles",
    "",
    "🔴 PRIORITY 2 - Polyester Chips/Fiber: ₹3.2-3.6B import → 75% substitution by FY28",
    "    Current: 40% imported | Suppliers: Indorama, IOCL Bhadrak (+1.2 MMTPA combined)",
    "    Cost Savings: $100-250/MT | Affected Segment: CLOTHTECH weaving, apparel",
    "",
    "🔴 PRIORITY 3 - Polyethylene: ₹5.6-6.0B import → 75% substitution by FY30",
    "    Current: 40% imported | Supplier: RIL O2C (+1.8 MMTPA)",
    "    Cost Savings: $50-150/MT | Affected Segment: PACKTECH films, laminates"
])

# ============================================================================
# SLIDE 7: PRODUCT-WISE HSN CODE MAPPING
# ============================================================================
hsn_data = [
    ["Product", "HSN Code", "Current Import (₹ Cr)", "Substitution %", "Companies to Target"],
    ["LLDPE (Film)", "3914", "800-900", "80%", "Nonwoven/film converters"],
    ["Polyester Chips", "3907", "1,600-1,800", "75%", "Weavers, integrated mills"],
    ["Polyester Filament", "5402", "600-700", "70%", "Knit mills, fabric makers"],
    ["Specialty Dyes", "3204", "300-400", "60%", "Dye manufacturers"],
    ["Aromatic Hydrocarbons", "2903", "600-700", "70%", "Chemical intermediates"]
]
add_table_slide(prs, "HSN Code Mapping - High Priority Products", hsn_data)

# ============================================================================
# SLIDE 8: SEGMENT-WISE MARGIN RECOVERY
# ============================================================================
seg_data = [
    ["Segment", "Current Margin", "FY30 Target", "Profit Swing (₹ Cr)", "Key Driver"],
    ["HOMETECH", "89%", "89%+", "0 (LOCKED)", "Scale + dyes integration"],
    ["CLOTHTECH", "+15.8%", "+66.7%", "300-400", "Polyester cost parity"],
    ["INDUTECH", "-80.9%", "+20%", "210-280", "LLDPE substitution (CRITICAL)"],
    ["MOBILTECH", "-50%", "+25%", "150-200", "Specialty polymers (BPCL)"],
    ["PACKTECH", "89%", "89%+", "0 (PROTECTED)", "RIL PE supply"],
    ["TOTAL UPLIFT", "-", "-", "810-1,130", "FY28-30 realization"]
]
add_table_slide(prs, "Segment-wise Margin Recovery Opportunity", seg_data)

# ============================================================================
# SLIDE 9: NEW COMPANY TARGETING FRAMEWORK
# ============================================================================
add_content_slide(prs, "New Company Targeting Framework", [
    "✓ LLDPE Entry Profile (10-15 companies): INDUTECH converters with ₹300-500 KTPA imports",
    "    → Capex: ₹400-600 Cr | ROI: 4-6 years | Cost savings: $150-200/MT × 300 KTPA = $45-60M/yr",
    "",
    "✓ Polyester Entry Profile (15-20 companies): Regional weavers with ₹800-1,000 KTPA imports",
    "    → Capex: ₹500-800 Cr | ROI: 4-6 years | Cost savings: $100-200/MT × 600 KTPA = $60-120M/yr",
    "",
    "✓ PE/PP Entry Profile (5-10 companies): Film extruders with ₹400-800 KTPA consumption",
    "    → Capex: ₹200-400 Cr | ROI: 4-6 years | Cost savings: $50-150/MT × 400 KTPA = $20-60M/yr",
    "",
    "✓ Specialty Chemical Entry (5-10 companies): Dye/pigment makers for ecosystem integration",
    "    → Capex: ₹150-350 Cr | ROI: 3-5 years | Margin uplift: 30-35% specialty focus"
])

# ============================================================================
# SLIDE 10: CAPEX TIMELINE & PHASING
# ============================================================================
add_content_slide(prs, "Capex Timeline: FY26-FY30 Phasing", [
    "📌 Phase 1 (FY26): Contract & Land Allocation",
    "    Spend: ₹15-20K Cr | BPCL LLDPE/Indorama supply contracts (DEADLINE: Mar 2026)",
    "    Action: Secure land, finalize equipment specs, board approvals",
    "",
    "📌 Phase 2 (FY27): Construction & Equipment Installation — PEAK YEAR",
    "    Spend: ₹50-60K Cr | RIL O2C, BPCL, Trident, Jain Cord, Indorama construction",
    "    Action: Equipment delivery, civil works, factory construction",
    "",
    "📌 Phase 3 (FY28): Production Ramp & First Shipments",
    "    Spend: ₹40-50K Cr | BPCL Bina LLDPE online (Q3 FY28), production scaling",
    "    Action: Commissioning, quality validation, customer ramp-up",
    "",
    "📌 Phase 4 (FY29-30): Full Capacity & Stabilization",
    "    Spend: ₹30-40K Cr | All projects operational, margin recovery peaks (FY28-30)"
])

# ============================================================================
# SLIDE 11: GEOGRAPHIC DISTRIBUTION
# ============================================================================
geo_data = [
    ["Region", "Capex (₹ Cr)", "Key Companies", "Strategic Advantage"],
    ["Coastal (PCPIR + Ports)", "138,000", "RIL, Indorama, IOCL, BPCL", "Port access, refining base"],
    ["Inland (Dhar, Rajasthan)", "100,000", "PM Mitra parks, textiles", "Cotton zones, cost advantage"],
    ["Emerging (Bihar, others)", "12,000", "Regional mills, MSME", "New industrial corridors"],
    ["TOTAL ECOSYSTEM", "250,000", "50-80 companies", "All-India fabric coverage"]
]
add_table_slide(prs, "Investment Distribution by Geographic Region", geo_data)

# ============================================================================
# SLIDE 12: CRITICAL SUCCESS FACTORS
# ============================================================================
add_content_slide(prs, "Critical Success Factors & Risk Mitigation", [
    "🔴 URGENT: Lock BPCL Bina LLDPE supply contracts by Mar 2026",
    "    → First-come-first-served allocation | Miss deadline = 1-2 year delay in margin recovery",
    "    → $150-200/MT savings only for early signers | Cost parity lost for late entrants",
    "",
    "🟡 Manage execution risk in PM Mitra parks (historically 50-70% completion rate)",
    "    → Anchor tenants (Trident, Jain Cord, Arvind) de-risked with land allotment",
    "    → Smaller players require government subsidy coordination (40% machinery, 20-30% power)",
    "",
    "🟢 RIL O2C & BPCL capex on-track (environmental clearance progressing)",
    "    → Feedstock supply guaranteed by FY27-28 for downstream conversion capex",
    "    → All downstream players should stage capex to align with feedstock online dates"
])

# ============================================================================
# SLIDE 13: GO-TO-MARKET STRATEGY FOR NEW ENTRANTS
# ============================================================================
add_content_slide(prs, "Go-to-Market Strategy for New Entrants", [
    "PHASE 1 (FY26 - Q1-Q2): Contracting & Infrastructure",
    "    1. Identify BPCL/Indorama supply contracts (deadline Jun 2026)",
    "    2. Secure land (10-25 acres) within 500 km of feedstock supplier",
    "    3. Prepare project reports, secure board/financing approval",
    "",
    "PHASE 2 (FY27-28): Construction & Production Ramp",
    "    1. Equipment procurement (12-18 month lead time)",
    "    2. Pilot production & buyer quality validation",
    "    3. Scale to 50-70% design capacity by FY28-Q3",
    "",
    "PHASE 3 (FY28-30): Market Capture & Margin Realization",
    "    1. Secure long-term customer contracts (3-5 year supply locks)",
    "    2. Achieve ₹150-300M annual profit per ₹500-800 Cr capex",
    "    3. Plan Phase 2 expansion (adjacent products, markets)"
])

# ============================================================================
# SLIDE 14: INVESTMENT SCORING MATRIX
# ============================================================================
score_data = [
    ["Company", "Capex", "Segment", "Timing", "Risk", "Score"],
    ["Trident", "4.9K Cr", "HOMETECH", "FY28 ✓", "Low", "9/10"],
    ["Jain Cord", "2.5K Cr", "INDUTECH", "FY28 ✓", "Medium", "8.5/10"],
    ["Arvind Mills", "1.0K Cr", "CLOTHTECH", "FY28 ✓", "Medium", "8/10"],
    ["AB Cotspin", "1.3K Cr", "Yarn", "FY27 ✓", "Low-Med", "7.5/10"],
    ["RSVPM", "0.7K Cr", "MOBILTECH", "FY28 ✓", "Medium", "7.5/10"],
    ["RE&UP", "4.8K Cr", "Recycling", "FY26-28", "High", "6.5/10"]
]
add_table_slide(prs, "Investment Scoring Matrix (Score >8 = Core, 6-8 = Opportunistic)", score_data)

# ============================================================================
# SLIDE 15: TOP 10 NEW ENTRY COMPANIES TO TARGET
# ============================================================================
entry_data = [
    ["Rank", "Profile", "Segment", "Capex", "Opportunity (₹ Cr)"],
    ["1", "Tier-2 nonwoven producer (₹50-100 Cr rev)", "LLDPE", "400-600", "510-710"],
    ["2", "Regional weaving mill (₹100-200 Cr rev)", "Polyester", "500-800", "850-1,050"],
    ["3", "Film converter (₹50-100 Cr rev)", "PE", "200-400", "200-300"],
    ["4", "Auto-textile supplier", "Specialty polymers", "300-500", "150-250"],
    ["5", "Chemical company (dye production)", "Specialty dyes", "150-350", "100-150"],
    ["6", "Geotextile producer", "PE/LLDPE blend", "200-400", "150-250"],
    ["7", "Apparel manufacturer", "Polyester fiber", "300-500", "100-200"],
    ["8", "Textile chemical firm", "Finishing agents", "150-300", "100-150"],
    ["9", "Pigment/paint company", "Pigments", "100-250", "50-100"],
    ["10", "Regional specialty cluster", "Multi-segment", "800-1,500", "500-1,000"]
]
add_table_slide(prs, "Top 10 New Company Entry Profiles by Segment & Opportunity", entry_data)

# ============================================================================
# SLIDE 16: KEY METRICS - BEFORE & AFTER
# ============================================================================
add_content_slide(prs, "Expected Impact by FY30 (2029-30)", [
    "📊 Import Substitution: ₹18,500-20,600 Cr annual imports → 70-80% domestically sourced",
    "💰 Annual Margin Recovery: ₹810-1,130 Crore profit uplift across ecosystem",
    "",
    "🏭 Domestic Feedstock Capacity Additions:",
    "   • Polyethylene: +1.8 MMTPA (RIL) | LLDPE: +1.0 MMTPA (BPCL)",
    "   • Polyester: +1.2 MMTPA (Indorama + IOCL) | Specialty: +0.9 MMTPA (BPCL AP)",
    "",
    "💼 Investment Catalysts:",
    "   • 50-80 new companies entering ecosystem | 300,000+ direct jobs created",
    "   • ₹250,000 Cr capex deployed over 5 years | Export competitiveness +40% by FY30",
    "",
    "⚡ Cost Advantage: Domestic polyester, LLDPE, PE cost parity achieved by FY28-30"
])

# ============================================================================
# SLIDE 17: IMMEDIATE ACTION ITEMS (NEXT 90 DAYS)
# ============================================================================
add_content_slide(prs, "Action Items: Next 90 Days (Jul-Sep 2026)", [
    "🎯 BPCL LLDPE Supply Contracts (DEADLINE: Mar 2026 — ALREADY PASSED, execute immediately)",
    "   Action: Engage procurement teams, negotiate 3-5yr offtake agreements",
    "   Impact: $150-200/MT cost savings locked for FY27-28 production ramp",
    "",
    "🎯 Indorama Polyester Long-term Pricing",
    "   Action: Secure allocation from Paradeep/Kalyan expansion (first-come-first-served)",
    "   Impact: ₹300-400M annual margin recovery for CLOTHTECH ecosystem",
    "",
    "🎯 PM Mitra Land Allotment Confirmation",
    "   Action: 90% of 91 companies should confirm land allocation by Jun 2026",
    "   Impact: Enables FY27 capex spend acceleration",
    "",
    "🎯 Identify & Shortlist New Entry Companies",
    "   Action: Outreach to 20-30 nonwoven, weaving, film converter companies",
    "   Impact: 50-80 company ecosystem lock-in for FY26-30 capex pipeline"
])

# ============================================================================
# SLIDE 18: CLOSING - KEY TAKEAWAYS
# ============================================================================
add_content_slide(prs, "Key Takeaways", [
    "1️⃣  ₹250,000+ Crore ecosystem capex opportunity over FY26-30 (72% in Tier-1 companies)",
    "",
    "2️⃣  Import substitution potential: ₹4,500-5,500 Cr annually by FY30 (70-80% penetration)",
    "",
    "3️⃣  Margin recovery concentrated in 3 segments: INDUTECH (₹210-280M), CLOTHTECH (₹300-400M), MOBILTECH (₹150-200M)",
    "",
    "4️⃣  Critical deadline: Mar 2026 for BPCL Bina LLDPE contracts → Every month delay = $30-50M annual opportunity loss",
    "",
    "5️⃣  New company targeting: 50-80 companies in LLDPE, polyester, PE, specialty chemicals ready for entry",
    "",
    "6️⃣  Timeline: FY27 is peak capex year (₹50-60K Cr spend) → Production ramp FY27-28 → Full margin realization FY28-30"
])

# ============================================================================
# Save presentation
# ============================================================================
output_path = "/Users/umashankar/omc-retail-profitability-model/textile_data/research/Textile_Chemical_Investment_Analysis.pptx"
prs.save(output_path)
print(f"✅ PowerPoint presentation created successfully!")
print(f"📄 File: {output_path}")
print(f"📊 Slides: 18 (Title + 17 content slides)")
print(f"\nPresentation includes:")
print("  • Executive summary & investment landscape")
print("  • Master investment list (Tier 1-5 companies)")
print("  • Import substitution opportunities by product & HSN code")
print("  • Segment-wise margin recovery analysis")
print("  • New company targeting framework")
print("  • Go-to-market strategy & timeline")
print("  • Critical success factors & immediate actions")
